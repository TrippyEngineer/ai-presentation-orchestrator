"""
slide_reader.py — Extracts text content from each slide in the PPTX.

Configure PPTX_FILE in .env or pass path directly to extract_slide_content().
Configure SKIP_CONTENT_SLIDES and NOISE_STRINGS in .env to match your deck.
"""

import os
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

PPTX_FILE = os.getenv("PPTX_FILE", "presentation.pptx")

# Slide numbers that are purely visual — no useful text to extract (1-indexed).
# Override via SKIP_SLIDES env var: e.g. SKIP_SLIDES=2,16
_skip_raw = os.getenv("SKIP_SLIDES", "")
SKIP_CONTENT_SLIDES = {int(x) for x in _skip_raw.split(",") if x.strip().isdigit()}

# Text strings to ignore when extracting content (e.g. footer/watermark text).
# Override via NOISE_STRINGS env var (comma-separated).
_noise_raw = os.getenv("NOISE_STRINGS", "")
NOISE_STRINGS = {s.strip() for s in _noise_raw.split(",") if s.strip()}


def extract_slide_content(pptx_path: str = PPTX_FILE) -> list[dict]:
    """
    Extract text content from each slide.
    Returns list of dicts: [{slide_number, content, has_content}]
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(
            f"PPTX not found at: {pptx_path}\n"
            f"Set PPTX_FILE in .env or place presentation.pptx in this folder."
        )

    prs = Presentation(pptx_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        text_parts = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            if text in NOISE_STRINGS:
                continue
            text_parts.append(text)

        content = "\n".join(text_parts)
        has_content = bool(content.strip()) and slide_num not in SKIP_CONTENT_SLIDES

        slides_content.append({
            "slide_number": slide_num,
            "content": content,
            "has_content": has_content,
        })

    return slides_content


if __name__ == "__main__":
    slides = extract_slide_content()
    print(f"\nFound {len(slides)} slides:\n")
    for slide in slides:
        status = "✅" if slide["has_content"] else "⬛"
        words = len(slide["content"].split())
        print(f"  {status} Slide {slide['slide_number']:2d}: {words:4d} words")
        if slide["content"]:
            preview = slide["content"].replace("\n", " ")[:80]
            print(f"     {preview}...")
