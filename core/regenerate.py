"""
regenerate.py — Interactive script + audio regenerator for AI Workshop.

CHANGES v2:
  - Slide content read DYNAMICALLY from PPTX file (no hardcoded dict)
  - Slide structure automatically reflects PPT changes (deletions, reordering)
  - Cache is validated against current PPTX structure and flagged if stale
  - Logging via logger.py
  - Voice updated to en-IN-PrabhatNeural (Indian English male)
  - Rate updated to +8% (slightly faster)

Run: python regenerate.py
"""

import os, sys, json, time, asyncio, subprocess, platform
from pathlib import Path
from dotenv import load_dotenv

from logger import get_logger, log_slide_event

load_dotenv()
log = get_logger("regenerate")

CACHE_DIR     = Path("cache")
SCRIPTS_PATH  = CACHE_DIR / "scripts.json"
MANIFEST_PATH = CACHE_DIR / "manifest.json"
PPTX_FILE     = "AI_Workshop_Presentation.pptx"
CACHE_DIR.mkdir(exist_ok=True)

VOICE = "en-IN-PrabhatNeural"
RATE  = "+8%"


# ══════════════════════════════════════════════════════════════════════
# DYNAMIC SLIDE CONTENT FROM PPTX
# ══════════════════════════════════════════════════════════════════════

def read_slide_content_from_pptx(pptx_path: str = PPTX_FILE) -> dict[int, str]:
    """
    Read actual slide content from the PPTX file.
    Returns {slide_number: text_content}.
    This ensures regenerate.py always reflects the CURRENT deck structure.
    """
    from pptx import Presentation

    if not Path(pptx_path).exists():
        log.error(f"PPTX not found: {pptx_path}")
        raise FileNotFoundError(
            f"PPTX not found at: {pptx_path}\n"
            f"Place {pptx_path} in the same folder as this script."
        )

    prs = Presentation(pptx_path)
    slide_content = {}
    skip_noise = {"AI Workshop — Confidential", "AI Workshop"}

    for i, slide in enumerate(prs.slides):
        slide_num  = i + 1
        text_parts = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text or text in skip_noise:
                continue
            text_parts.append(text)
        slide_content[slide_num] = "\n".join(text_parts)

    total = len(slide_content)
    log.info(f"Read {total} slides from {pptx_path}")
    return slide_content


def detect_demo_slides(slide_content: dict[int, str]) -> dict[int, str]:
    """
    Auto-detect demo slides by searching for 'LIVE DEMO' or 'demo' trigger text.
    Returns {slide_num: demo_type}
    """
    demo_map = {}
    for num, content in slide_content.items():
        lower = content.lower()
        if "live demo 1" in lower or ("demo" in lower and "email" in lower):
            demo_map[num] = "email"
        elif "live demo 2" in lower or ("demo" in lower and "meeting" in lower):
            demo_map[num] = "meeting"
        elif "live demo 3" in lower or ("agentic" in lower and "demo" in lower):
            demo_map[num] = "research"
    log.info(f"Auto-detected demo slides: {demo_map}")
    return demo_map


def validate_cache_vs_pptx():
    """
    Check if cached scripts/manifest match the current PPTX structure.
    Warns if slides have been added, removed, or content changed.
    """
    print("\n  [ CACHE VALIDATION ]")
    try:
        current_slides = read_slide_content_from_pptx()
        current_count  = len(current_slides)
    except FileNotFoundError as e:
        print(f"  ⚠️  Cannot validate: {e}")
        return

    if SCRIPTS_PATH.exists():
        cached = json.loads(SCRIPTS_PATH.read_text())
        cached_count = len(cached)
        if cached_count != current_count:
            print(f"  ⚠️  SLIDE COUNT MISMATCH: cache has {cached_count}, PPTX now has {current_count}")
            log.warning(f"Cache/PPTX mismatch: cached={cached_count}, pptx={current_count}")
            print(f"     Some cached audio may be for slides that no longer exist.")
            print(f"     Recommendation: regenerate affected slides or run python pre_generate.py --force")

            stale_keys = [k for k in cached if int(k) > current_count]
            if stale_keys:
                print(f"  ⚠️  Stale cached slides (beyond current deck): {stale_keys}")
        else:
            print(f"  ✅ Cache matches PPTX: {current_count} slides")
            log.info(f"Cache validation OK: {current_count} slides")
    else:
        print(f"  ℹ️  No cache yet. PPTX has {current_count} slides.")

    print()


# ══════════════════════════════════════════════════════════════════════
# STEP 1: GUIDED PROMPT INTERVIEW
# ══════════════════════════════════════════════════════════════════════

def run_prompt_interview(slide_nums: list[int]) -> dict:
    print("\n" + "─"*65)
    print("  SCRIPT CUSTOMISATION INTERVIEW")
    print("  Answer each question. Press ENTER to keep the default.")
    print("─"*65 + "\n")

    prefs = {}

    print("  Q1. LANGUAGE LEVEL")
    print("      1 = Plain English only    2 = Mostly plain    3 = Balanced    4 = Technical")
    ans = input("      Your choice (1-4) [default: 2]: ").strip() or "2"
    levels = {
        "1": "Use plain English only. Avoid ALL technical jargon. Explain every term simply.",
        "2": "Use mostly plain language. Brief explanations for technical terms. Audience uses basic software but is not technical.",
        "3": "Balanced. Audience is comfortable with automation and AI tools concepts.",
        "4": "Technical language fine. Audience understands APIs, automation, agentic AI.",
    }
    prefs["language"] = levels.get(ans, levels["2"])

    print("\n  Q2. TONE")
    print("      1 = Calm/authoritative    2 = Warm/collegial    3 = Energetic    4 = Data-driven")
    ans = input("      Your choice (1-4) [default: 2]: ").strip() or "2"
    tones = {
        "1": "Calm, measured, authoritative. No hype. Confident expertise.",
        "2": "Warm and collegial. Like a trusted colleague who cares about the team's success.",
        "3": "Energetic and inspiring. Build excitement. This is a turning point.",
        "4": "Serious and data-led. Lead with facts. Trust the data to make the case.",
    }
    prefs["tone"] = tones.get(ans, tones["2"])

    print("\n  Q3. SCRIPT LENGTH PER SLIDE")
    print("      1 = Short (70-100 words)    2 = Medium (110-150 words)    3 = Long (150-190 words)")
    ans = input("      Your choice (1-3) [default: 2]: ").strip() or "2"
    lengths = {
        "1": "70-100 words per slide. Punchy. Leave space for visuals.",
        "2": "110-150 words per slide. Standard presentation pace. Cover key points naturally.",
        "3": "150-190 words per slide. Thorough. More context and examples.",
    }
    prefs["length"] = lengths.get(ans, lengths["2"])

    print("\n  Q4. MISSION EMPHASIS")
    print("      1 = Minimal    2 = Moderate    3 = Strong (connect every capability to health impact)")
    ans = input("      Your choice (1-3) [default: 2]: ").strip() or "2"
    missions = {
        "1": "Focus on the technology. Mention  work only when directly relevant.",
        "2": "Connect AI capabilities back to mission at natural moments.",
        "3": "Consistently ground every capability in health outcomes and mission.",
    }
    prefs["mission"] = missions.get(ans, missions["2"])

    print("\n  Q5. CONCRETE EXAMPLES")
    print("      1 = Generic    2 = Global health examples    3 = Examples-first approach")
    ans = input("      Your choice (1-3) [default: 2]: ").strip() or "2"
    examples = {
        "1": "Keep examples generic. Do not name specific diseases, regions, or programs.",
        "2": "Use specific examples: malaria interventions, maternal health, Sahel region, USAID briefs, Kenya field teams.",
        "3": "Build every slide around a vivid global health example. Example first, then technology.",
    }
    prefs["examples"] = examples.get(ans, examples["2"])

    print("\n  Q6. AUDIENCE FAMILIARITY WITH AI")
    print("      1 = Beginners    2 = Casual users    3 = Engaged users")
    ans = input("      Your choice (1-3) [default: 2]: ").strip() or "2"
    familiarity = {
        "1": "Assume beginners. Define every AI concept as you go.",
        "2": "Assume occasional ChatGPT users. Show them what's possible at scale.",
        "3": "Assume regular AI users. Skip basics. Focus on infrastructure-level possibilities.",
    }
    prefs["familiarity"] = familiarity.get(ans, familiarity["2"])

    print(f"\n  Q7. PER-SLIDE SPECIFIC INSTRUCTIONS (optional)")
    print(f"      Slides being regenerated: {slide_nums}")
    print("      Format: <slide_number>: <instruction>   Type DONE when finished.\n")
    slide_instructions = {}
    while True:
        line = input("      Instruction (or DONE): ").strip()
        if line.upper() == "DONE" or line == "":
            break
        if ":" not in line:
            print("      Use format:  <slide number>: <instruction>")
            continue
        parts = line.split(":", 1)
        try:
            num   = int(parts[0].strip())
            instr = parts[1].strip()
            slide_instructions[num] = instr
            print(f"      ✓ Slide {num} instruction saved")
        except ValueError:
            print("      Invalid slide number")
    prefs["slide_instructions"] = slide_instructions

    confirm = input("\n  Proceed with these settings? (y/n): ").strip().lower()
    if confirm != "y":
        print("  Cancelled.\n")
        sys.exit(0)
    return prefs


# ══════════════════════════════════════════════════════════════════════
# STEP 2: GENERATE SCRIPT VIA CLAUDE
# ══════════════════════════════════════════════════════════════════════

def build_prompt(slide_num: int, slide_content: str, total_slides: int, prefs: dict,
                 demo_map: dict) -> str:
    is_demo   = slide_num in demo_map
    demo_type = demo_map.get(slide_num, "")
    slide_instruction = prefs.get("slide_instructions", {}).get(slide_num, "")

    length_str = prefs.get("length", "")
    if "70-100"  in length_str: word_range = "70–100"
    elif "150-190" in length_str: word_range = "150–190"
    else: word_range = "110–150"

    if is_demo:         word_range = "65–90"
    if slide_num == total_slides: word_range = "60–90"   # closing slide

    demo_ctx = {
        "email":    "About to trigger email workflow live. Tell audience to watch — donor email arriving. Not touching anything. Short. Tense.",
        "meeting":  "About to trigger meeting pipeline live. One transcript in, structured outputs out. Not touching anything. Short. Crisp.",
        "research": "Most impressive demo — agentic AI. Previous demos had fixed paths. This one decides its own. Watch it plan, search, evaluate, choose strategy. Build the distinction. Then hand over.",
    }.get(demo_type, "")

    custom_block = f"\nSPECIFIC INSTRUCTION: {slide_instruction}" if slide_instruction else ""
    demo_block   = f"\nDEMO CONTEXT: {demo_ctx}" if demo_ctx else ""

    return f"""You are [YOUR NAME], digital health and AI expert presenting to colleagues in India.

VOICE SETTINGS:
- Language: {prefs['language']}
- Tone: {prefs['tone']}
- Mission emphasis: {prefs['mission']}
- Examples: {prefs['examples']}
- Audience: {prefs['familiarity']}

RULES:
- Never say "This slide shows", "As you can see", "In conclusion", "Let me now", "Moving on"
- Open with a scene, question, or vivid example
- Use "you" / "your team" at least once to pull the audience in
- Short punchy lines mixed with fuller explanations for rhythm
- Output ONLY speech text — no labels, no markdown, no stage directions
- Word count: {word_range} words exactly

SLIDE {slide_num} of {total_slides} — CURRENT CONTENT:
{slide_content.strip() if slide_content.strip() else "(Visual/video slide — minimal speech)"}
{demo_block}
{custom_block}

Write exactly what [YOUR NAME] says when this slide appears on screen."""


def generate_script(slide_num: int, slide_content: str, total_slides: int,
                    prefs: dict, demo_map: dict) -> str:
    """Call Claude to generate one slide's script."""
    try:
        from anthropic import Anthropic
        qa_client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    except ImportError:
        log.error("anthropic not installed")
        print("  ❌ anthropic not installed. Run: pip install anthropic")
        return ""

    prompt = build_prompt(slide_num, slide_content, total_slides, prefs, demo_map)
    log.info(f"Generating script for slide {slide_num}")

    try:
        response = qa_client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=600,
            messages=[{"role": "user", "content": prompt}]
        )
        speech = response.content[0].text.strip()
        log.info(f"Slide {slide_num}: generated {len(speech.split())} words")
        return speech
    except Exception as e:
        log.error(f"Claude API error on slide {slide_num}: {e}", exc_info=True)
        print(f"  ❌ Claude API error: {e}")
        return ""


# ══════════════════════════════════════════════════════════════════════
# STEP 3: AUDIO SYNTHESIS
# ══════════════════════════════════════════════════════════════════════

async def _synth_async(text: str, path: str) -> bool:
    try:
        import edge_tts
        comm = edge_tts.Communicate(text, VOICE, rate=RATE)
        await comm.save(path)
        return True
    except ImportError:
        log.error("edge-tts not installed")
        print("  ❌ edge-tts not installed. Run: pip install edge-tts")
        return False
    except Exception as e:
        log.error(f"TTS error for {path}: {e}")
        print(f"  ❌ TTS error: {e}")
        return False


def synthesise(text: str, path: str) -> bool:
    if not text.strip():
        return False
    return asyncio.run(_synth_async(text, path))


def play_audio(path: str):
    if not Path(path).exists():
        print(f"  No audio at {path}")
        return
    played = False
    try:
        import pygame
        if not pygame.mixer.get_init():
            pygame.mixer.pre_init(44100, -16, 2, 2048)
            pygame.mixer.init()
        pygame.mixer.music.stop()
        pygame.mixer.music.load(path)
        pygame.mixer.music.play()
        t = time.time()
        while pygame.mixer.music.get_busy() and time.time() - t < 120:
            time.sleep(0.1)
        pygame.mixer.music.stop()
        played = True
    except Exception:
        pass
    if not played:
        try:
            subprocess.run(
                ["ffplay", "-nodisp", "-autoexit", "-loglevel", "quiet", path],
                timeout=120, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
            )
            played = True
        except Exception:
            pass
    if not played and platform.system() == "Windows":
        try:
            import ctypes
            winmm = ctypes.windll.winmm
            alias = "regenplay"
            p = str(Path(path).resolve())
            winmm.mciSendStringW(f'open "{p}" type mpegvideo alias {alias}', None, 0, None)
            winmm.mciSendStringW(f"play {alias} wait", None, 0, None)
            winmm.mciSendStringW(f"close {alias}", None, 0, None)
        except Exception:
            pass


def get_duration(path: str) -> float:
    try:
        r = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", path],
            capture_output=True, text=True, timeout=10
        )
        if r.returncode == 0 and r.stdout.strip():
            return float(r.stdout.strip())
    except Exception:
        pass
    try:
        return max(Path(path).stat().st_size / 6000, 1.0)
    except Exception:
        return 5.0


# ══════════════════════════════════════════════════════════════════════
# STEP 4: SAVE + REBUILD MANIFEST
# ══════════════════════════════════════════════════════════════════════

def load_scripts() -> dict:
    if SCRIPTS_PATH.exists():
        return json.loads(SCRIPTS_PATH.read_text())
    return {}


def save_scripts(scripts: dict):
    SCRIPTS_PATH.write_text(json.dumps(scripts, indent=2))


def rebuild_manifest(scripts: dict, demo_map: dict):
    manifest = {}
    for num_str, text in scripts.items():
        num  = int(num_str)
        mp3  = CACHE_DIR / f"slide_{num:02d}.mp3"
        dur  = get_duration(str(mp3)) if mp3.exists() else 3.0
        manifest[num_str] = {
            "slide_num":  num,
            "audio_path": str(mp3) if mp3.exists() else None,
            "duration":   round(dur, 2),
            "speech":     text,
            "is_demo":    num in demo_map,
            "demo_type":  demo_map.get(num),
        }
    MANIFEST_PATH.write_text(json.dumps(manifest, indent=2))
    total_dur = sum(v["duration"] for v in manifest.values())
    log.info(f"Manifest rebuilt: {len(manifest)} slides, {total_dur/60:.1f}min")
    print(f"  Manifest saved — {len(manifest)} slides, {total_dur/60:.1f} min total")


# ══════════════════════════════════════════════════════════════════════
# MAIN FLOWS
# ══════════════════════════════════════════════════════════════════════

def flow_specific_slides(slide_content: dict, demo_map: dict):
    total_slides = max(slide_content.keys()) if slide_content else 18

    print(f"\n  PPTX has {total_slides} slides. Demo slides detected: {list(demo_map.keys())}")
    print("  Which slides do you want to regenerate?")
    print("  Enter slide numbers (e.g. 3, 8, 16), a range (9-13), or 'all'\n")

    raw = input("  Slides: ").strip()
    if not raw:
        print("  Nothing entered.")
        return

    slide_nums = []
    if raw.lower() == "all":
        slide_nums = sorted(slide_content.keys())
    else:
        for part in raw.replace(" ", "").split(","):
            if "-" in part:
                try:
                    a, b = part.split("-")
                    slide_nums.extend(range(int(a), int(b) + 1))
                except ValueError:
                    print(f"  Skipping invalid range: {part}")
            else:
                try:
                    slide_nums.append(int(part))
                except ValueError:
                    print(f"  Skipping: {part}")

    slide_nums = sorted(set(n for n in slide_nums if n in slide_content))
    if not slide_nums:
        print("  No valid slide numbers found in current PPTX.")
        return

    print(f"\n  Regenerating slides: {slide_nums}")
    prefs   = run_prompt_interview(slide_nums)
    scripts = load_scripts()

    for slide_num in slide_nums:
        print(f"\n  ══ SLIDE {slide_num}/{total_slides} ══════════════════════════")
        content = slide_content.get(slide_num, "")
        preview = content[:100].replace("\n", " ")
        print(f"  PPTX content: {preview}{'...' if len(content) > 100 else ''}")

        # Fixed content for video slides (no text on screen)
        if not content.strip():
            speech = "Good morning. Before I say anything else — press one button with me."
            print("  (Visual/video slide — using short bridge)")
        else:
            print("  Generating script via Claude...", end=" ", flush=True)
            speech = generate_script(slide_num, content, total_slides, prefs, demo_map)
            if not speech:
                print("  FAILED — skipping.")
                log.error(f"Slide {slide_num}: generation failed, skipping")
                continue

        words = len(speech.split())
        print(f"done ({words} words)")
        print(f"\n  ┌─ SCRIPT ────────────────────────────────────────────")
        for line in [speech[i:i+65] for i in range(0, len(speech), 65)]:
            print(f"  │ {line}")
        print(f"  └─────────────────────────────────────────────────────\n")

        action = input("  Accept (a) / Edit (e) / Regenerate (r) / Skip (s): ").strip().lower()

        if action == "s":
            print("  Skipped.")
            continue
        elif action == "r":
            extra = input("  Extra instruction for retry (or ENTER): ").strip()
            if extra:
                prefs.setdefault("slide_instructions", {})[slide_num] = extra
            speech = generate_script(slide_num, content, total_slides, prefs, demo_map)
            if not speech:
                print("  Regeneration failed. Skipping.")
                continue
            print(f"\n  New ({len(speech.split())} words): {speech}\n")
            action = input("  Accept (a) / Edit (e) / Skip (s): ").strip().lower()
            if action == "s":
                continue
            elif action == "e":
                lines = []
                print("  Paste edited script, press Enter twice when done:\n")
                while True:
                    l = input("  ")
                    if l == "" and lines and lines[-1] == "":
                        break
                    lines.append(l)
                speech = "\n".join(l for l in lines if l).strip()
        elif action == "e":
            lines = []
            print("  Paste edited version, press Enter twice when done:\n")
            while True:
                l = input("  ")
                if l == "" and lines and lines[-1] == "":
                    break
                lines.append(l)
            speech = "\n".join(l for l in lines if l).strip()

        scripts[str(slide_num)] = speech
        save_scripts(scripts)
        log_slide_event(slide_num, "SCRIPT_SAVED", f"{len(speech.split())} words")
        print(f"  ✓ Script saved")

        mp3_path = str(CACHE_DIR / f"slide_{slide_num:02d}.mp3")
        if Path(mp3_path).exists():
            Path(mp3_path).unlink()
            log.debug(f"Deleted old audio: {mp3_path}")

        print("  Synthesising audio (Indian English, {RATE})...", end=" ", flush=True)
        ok = synthesise(speech, mp3_path)
        if ok:
            dur = get_duration(mp3_path)
            print(f"done ({dur:.1f}s)")
            play_it = input("  Play audio to verify? (y/n): ").strip().lower()
            if play_it == "y":
                play_audio(mp3_path)
        else:
            print("  FAILED — check internet and edge-tts")

    print("\n  Rebuilding manifest...")
    rebuild_manifest(scripts, demo_map)
    print("\n  ✓ Done. Run python orchestrator.py to present with updated audio.")


def flow_preview_scripts():
    scripts = load_scripts()
    if not scripts:
        print("\n  No scripts cached yet. Run python pre_generate.py first.\n")
        return
    print(f"\n  {'─'*65}")
    print(f"  {'SLIDE':<8} {'WORDS':<7} PREVIEW")
    print(f"  {'─'*65}")
    for k in sorted(scripts.keys(), key=lambda x: int(x)):
        text    = scripts[k]
        words   = len(text.split()) if text else 0
        preview = text[:70].replace("\n", " ") + ("..." if len(text) > 70 else "")
        print(f"  {int(k):<8} {words:<7} {preview}")
    total = sum(len(t.split()) for t in scripts.values())
    print(f"  {'─'*65}")
    print(f"  Total: {total} words across {len(scripts)} slides\n")
    view = input("  View full script for a slide? Enter number or ENTER to skip: ").strip()
    if view:
        try:
            k = str(int(view))
            if k in scripts:
                print(f"\n  ─── SLIDE {view} ───────────────────")
                print(f"  {scripts[k]}\n")
        except ValueError:
            pass


def flow_play_slide():
    raw = input("\n  Play which slide? Enter number: ").strip()
    try:
        num = int(raw)
        mp3 = CACHE_DIR / f"slide_{num:02d}.mp3"
        if mp3.exists():
            print(f"  Playing slide {num}...")
            play_audio(str(mp3))
            print("  Done.")
        else:
            print(f"  No audio for slide {num}. Run pre_generate.py or regenerate it first.")
    except ValueError:
        print("  Invalid number.")


# ══════════════════════════════════════════════════════════════════════
# ENTRY POINT
# ══════════════════════════════════════════════════════════════════════

def main():
    print("\n" + "═"*65)
    print("  PATH AI WORKSHOP — SCRIPT & AUDIO REGENERATOR v2")
    print("═"*65)
    log.info("regenerate.py started")

    if not os.getenv("ANTHROPIC_API_KEY"):
        print("\n  ❌ ANTHROPIC_API_KEY not set. Add it to .env and re-run.\n")
        sys.exit(1)

    # Load current PPTX structure dynamically
    print("\n  Reading current PPTX structure...")
    try:
        slide_content = read_slide_content_from_pptx()
        demo_map      = detect_demo_slides(slide_content)
        print(f"  ✅ {len(slide_content)} slides found in {PPTX_FILE}")
        print(f"  Demo slides detected: {dict(sorted(demo_map.items()))}")
    except FileNotFoundError as e:
        print(f"  ❌ {e}")
        sys.exit(1)

    # Validate cache against current deck
    validate_cache_vs_pptx()

    scripts     = load_scripts()
    audio_files = list(CACHE_DIR.glob("slide_*.mp3"))
    print(f"  Cached scripts : {len(scripts)} slides")
    print(f"  Audio files    : {len(audio_files)} files in cache/")

    while True:
        print("\n" + "─"*65)
        print("  MENU")
        print()
        print("  1. Regenerate specific slides (guided interview)")
        print("  2. Preview all current scripts + word counts")
        print("  3. Play audio for a specific slide")
        print("  4. Exit")
        print()
        choice = input("  Enter 1, 2, 3, or 4: ").strip()
        if choice == "1":
            flow_specific_slides(slide_content, demo_map)
        elif choice == "2":
            flow_preview_scripts()
        elif choice == "3":
            flow_play_slide()
        elif choice == "4":
            print("\n  Exiting. Run python orchestrator.py to present.\n")
            log.info("regenerate.py exited by user")
            break
        else:
            print("  Please enter 1, 2, 3, or 4.")


if __name__ == "__main__":
    main()
